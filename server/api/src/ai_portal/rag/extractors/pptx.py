"""PPTX extractor via :mod:`python-pptx`."""
from __future__ import annotations

from io import BytesIO
from typing import Any

from ai_portal.rag.extractors.protocol import (
    Block,
    ExtractedDocument,
    HeadingBlock,
    ParagraphBlock,
)


class PptxExtractor:
    name = "pptx"
    mime_types = {
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.ms-powerpoint",
    }

    def supports(self, mime: str) -> bool:
        return mime in self.mime_types

    async def extract(self, data: bytes, meta: dict[str, Any]) -> ExtractedDocument:
        from pptx import Presentation  # type: ignore

        prs = Presentation(BytesIO(data))
        blocks: list[Block] = []
        text_parts: list[str] = []
        for idx, slide in enumerate(prs.slides, start=1):
            title = None
            try:
                if slide.shapes.title is not None:
                    title = (slide.shapes.title.text or "").strip()
            except Exception:
                title = None
            if title:
                blocks.append(HeadingBlock(text=title, level=2, page=idx))
                text_parts.append(title)
            for shape in slide.shapes:
                if not getattr(shape, "has_text_frame", False):
                    continue
                for para in shape.text_frame.paragraphs:
                    txt = (para.text or "").strip()
                    if not txt or txt == title:
                        continue
                    blocks.append(ParagraphBlock(text=txt, page=idx))
                    text_parts.append(txt)
            # Slide notes appended.
            try:
                nf = slide.notes_slide.notes_text_frame if slide.has_notes_slide else None
                if nf is not None:
                    notes = (nf.text or "").strip()
                    if notes:
                        blocks.append(ParagraphBlock(text=notes, page=idx))
                        text_parts.append(notes)
            except Exception:
                pass
        return ExtractedDocument(
            text="\n".join(text_parts),
            blocks=blocks,
            meta={**meta, "slide_count": len(prs.slides)},
        )


__all__ = ["PptxExtractor"]
